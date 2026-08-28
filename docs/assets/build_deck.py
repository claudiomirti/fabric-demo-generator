"""Builds the executive narrative deck from docs/assets/narrative-slide.md.

The markdown file is the source of truth for the wording; this script lays it
out as a three-slide 16:9 deck. Re-run it after editing the narrative.

    pip install python-pptx
    python docs/assets/build_deck.py

Design notes:
  - Avoid emoji inside the Consolas comparison table. They render double-width
    but count as one character, which breaks the column alignment.
  - Check the rendered slides, not just the XML: text boxes overflow their
    background rectangles silently.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(0x1B, 0x1B, 0x1B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x00, 0x78, 0xD4)
RED = RGBColor(0xD1, 0x34, 0x38)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=14, bold=False, color=DARK, font_name='Segoe UI', space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = Pt(space_before)
    return p


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(slide, 0, 0, 13.333, 7.5, DARK)
add_text(slide, 1, 1.5, 11, 1.2, 'Fabric Data Demo Generator', size=44, bold=True, color=WHITE)
add_text(slide, 1, 2.8, 11, 0.8, 'Production-Shaped Demo Data for Microsoft Fabric — in Minutes',
         size=24, color=RGBColor(0x99, 0xCC, 0xFF))
tf = add_text(slide, 1, 4.2, 11, 1.5, 'Microsoft Fabric  +  Direct Lake  +  Data Agents',
              size=18, color=RGBColor(0xBB, 0xBB, 0xBB))
add_para(tf, 'Retail · Manufacturing · FSI · Healthcare · Life Science',
         size=16, color=RGBColor(0x88, 0x88, 0x88), space_before=12)

# Bottom line
add_rect(slide, 1, 5.55, 0.09, 0.6, BLUE)
add_text(slide, 1.28, 5.58, 11, 0.6,
         'From empty workspace to data-agent-ready model in five minutes.',
         size=22, bold=True, color=WHITE)
add_text(slide, 1, 6.5, 11, 0.4, 'github.com/claudiomirti/fabric-demo-generator',
         size=14, color=BLUE)

# ─── SLIDE 2: Problem + Why + How ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 0.9, RED)
add_text(slide, 0.5, 0.15, 12, 0.6, '🔴  THE PROBLEM', size=28, bold=True, color=WHITE)

# Problem box
add_rect(slide, 0.5, 1.2, 6, 2.8, LIGHT_BG)
tf = add_text(slide, 0.7, 1.3, 5.5, 2.5,
              'Every Fabric demo, POC, or data agent evaluation needs realistic data in a '
              'governed model — and building one by hand takes a day.',
              size=15, bold=True, color=DARK)
add_para(tf, '', size=8)
add_para(tf, '•  Toy CSVs — no star schema, nothing to show', size=13, color=GRAY, space_before=6)
add_para(tf, '•  Hand-built model — hours of clicking, easy to get wrong', size=13, color=GRAY, space_before=4)
add_para(tf, '•  Undescribed model — the data agent answers badly', size=13, color=GRAY, space_before=4)

# Why box
add_rect(slide, 6.8, 1.2, 6, 2.8, RGBColor(0xE8, 0xF4, 0xFD))
tf = add_text(slide, 7.0, 1.3, 5.5, 0.4, '💡  THE WHY', size=20, bold=True, color=BLUE)
add_para(tf, '', size=6)
add_para(tf, '"You can\'t show what a data agent can do on an empty workspace."',
         size=13, color=GRAY, space_before=8)
add_para(tf, '', size=6)
add_para(tf, '✓  Looks like a real business — a star schema', size=13, bold=True, color=DARK, space_before=6)
add_para(tf, '✓  Direct Lake on OneLake — the modern default', size=13, bold=True, color=DARK, space_before=4)
add_para(tf, '✓  Descriptions on every table, column, measure', size=13, bold=True, color=DARK, space_before=4)
add_para(tf, '✓  Repeatable — any workspace, in minutes', size=13, bold=True, color=DARK, space_before=4)

# ─── How section ───
add_rect(slide, 0.5, 4.3, 12.3, 0.5, BLUE)
add_text(slide, 0.7, 4.32, 11, 0.4, '⚙️  THE HOW  —  One App. Five Industries. Two Modes.',
         size=18, bold=True, color=WHITE)

for i, (title, sub, items) in enumerate([
    ('1 — Generate', 'Local Python', ['Fact + 3 dimensions', 'Up to 10,000 rows']),
    ('2 — Provision', 'Lakehouse + Delta', ['CSV → OneLake', 'Load Table API', 'Layout auto-detected']),
    ('3 — Model', 'Direct Lake + Ontology', ['Described model', 'DAX measures', 'Framed & ready']),
]):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 5.0, 3.8, 2.0, RGBColor(0xF0, 0xF0, 0xF0))
    add_text(slide, x + 0.1, 5.05, 3.6, 0.35, title, size=15, bold=True, color=BLUE)
    add_text(slide, x + 0.1, 5.35, 3.6, 0.25, sub, size=11, color=GRAY)
    tf2 = add_text(slide, x + 0.1, 5.65, 3.6, 1.2, '', size=12)
    for item in items:
        add_para(tf2, '•  ' + item, size=12, color=DARK, space_before=2)

    if i < 2:
        add_text(slide, x + 3.85, 5.7, 0.3, 0.4, '→', size=24, bold=True, color=BLUE)

# ─── SLIDE 3: What + Value ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 0.9, GREEN)
add_text(slide, 0.5, 0.15, 12, 0.6, '🎯  THE WHAT  +  🏆  THE VALUE', size=28, bold=True, color=WHITE)

# What section
add_rect(slide, 0.5, 1.2, 6, 3.2, RGBColor(0xE8, 0xF8, 0xE8))
tf = add_text(slide, 0.7, 1.3, 5.5, 0.4, 'One App, A Complete Demo Environment',
              size=18, bold=True, color=GREEN)
deliverables = [
    ('5 Industry Datasets', 'Retail, Mfg, FSI, Healthcare, Life Science'),
    ('Lakehouse + Delta Tables', 'CSV → OneLake → Load Table API'),
    ('Direct Lake Semantic Model', 'Relationships, DAX, descriptions'),
    ('Fabric Ontology', 'Entity types, properties, relationships'),
    ('Two Output Modes', 'CSV files only, or full provisioning'),
]
for title, desc in deliverables:
    add_para(tf, '', size=2)
    add_para(tf, '▸  ' + title, size=13, bold=True, color=DARK, space_before=3)
    add_para(tf, '    ' + desc, size=11, color=GRAY, space_before=1)

# Value comparison table
add_rect(slide, 6.8, 1.2, 6, 3.2, RGBColor(0xFD, 0xF4, 0xE8))
tf = add_text(slide, 7.0, 1.3, 5.5, 0.4, 'vs. Building the Demo by Hand',
              size=18, bold=True, color=ORANGE)
rows = [
    ('Realistic star schema, 5 industries', 'Hours', 'Minutes'),
    ('Direct Lake model built for you', 'No', 'Yes'),
    ('Descriptions for data agents', 'Skipped', 'Built in'),
    ('Large storage format + framing', 'Manual', 'Automatic'),
    ('Ontology item created', 'No', 'Yes'),
    ('Repeatable in any workspace', 'Manual', 'One click'),
]
add_para(tf, '', size=4)
add_para(tf, f'{"Capability":<34} {"By Hand":>9}  {"This":>9}',
         size=11, bold=True, color=DARK, space_before=6, font_name='Consolas')
for cap, f, t in rows:
    add_para(tf, f'{cap:<34} {f:>9}  {t:>9}',
             size=11, color=GRAY, space_before=2, font_name='Consolas')

# Bottom line
add_rect(slide, 0.5, 4.6, 12.3, 1.4, DARK)
tf = add_text(slide, 0.7, 4.7, 11.5, 0.5,
              '"Stop building demo data. Start demoing."',
              size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, '', size=8)
add_para(tf, 'Believable • Governed • Data-agent ready — provisioned straight into your own Fabric workspace',
         size=14, color=RGBColor(0xBB, 0xBB, 0xBB), space_before=8)

# Differentiators bar
add_rect(slide, 0.5, 6.3, 12.3, 0.9, LIGHT_BG)
for i, (dim, desc) in enumerate([
    ('Types Read from Delta', 'Never guessed from the CSV'),
    ('Lakehouse Layout Detected', 'Classic and schema-enabled'),
    ('Framed Before You Open It', 'First question returns an answer'),
]):
    x = 0.7 + i * 4.1
    add_text(slide, x, 6.35, 3.8, 0.3, dim, size=14, bold=True, color=BLUE)
    add_text(slide, x, 6.65, 3.8, 0.3, desc, size=12, color=GRAY)

out = Path(__file__).resolve().parent / 'Fabric-Demo-Generator-Narrative.pptx'
prs.save(out)
print(f'Saved: {out}')
